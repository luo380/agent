from pathlib import Path
try:
    from pypdf import PdfReader
except ImportError:  # pragma: no cover
    PdfReader = None

try:
    # DOCX 解析依赖：用于读取 Word 文档段落
    from docx import Document as DocxDocument
except ImportError:  # pragma: no cover
    DocxDocument = None

try:
    # XLSX 解析依赖：用于读取 Excel（xlsx）
    from openpyxl import load_workbook
except ImportError:  # pragma: no cover
    load_workbook = None

try:
    # XLS 解析依赖：用于读取老版 Excel（xls）
    import xlrd
except ImportError:  # pragma: no cover
    xlrd = None

try:
    # PPTX 解析依赖：用于读取 PowerPoint
    from pptx import Presentation
except ImportError:  # pragma: no cover
    Presentation = None




def normalize_text(text: str) -> str:
    """
    统一清洗文本格式，避免后续 chunking 时出现大量脏空行。

    处理策略：
    1. 去掉每行右侧空白
    2. 连续空行压缩为一个空行
    3. 最终去掉首尾空白
    """
    lines = [line.rstrip() for line in (text or "").splitlines()]
    cleaned: list[str] = []
    previous_blank = False

    for line in lines:
        is_blank = not line.strip()
        if is_blank:
            if not previous_blank:
                cleaned.append("")
            previous_blank = True
            continue

        cleaned.append(line.strip())
        previous_blank = False

    return "\n".join(cleaned).strip()


def _sheet_cell_to_text(value) -> str:
    """
    将 Excel 单元格的值安全转成字符串。

    这样做的目的是：
    - 数字、日期、布尔值都能统一落成文本
    - None 会被忽略，避免拼出无意义的 "None"
    """
    if value is None:
        return ""
    return str(value).strip()


def parse_txt(file_path: str) -> dict:
    """
    解析纯文本文件。

    纯文本没有天然页码和章节，所以：
    - pages 为空
    - sections 为空
    - full_text 存完整文本
    """
    text = Path(file_path).read_text(encoding="utf-8")
    return {
        "full_text": normalize_text(text),
        "pages": [],
        "sections": [],
        "metadata": {},
    }

def parse_md(file_path: str) -> dict:
    """
    解析 Markdown 文件。

    这里先做轻量实现：
    - 直接把 Markdown 作为文本读入
    - 不额外做标题树解析
    - 后续如需更细粒度 section，可再按 # / ## 扩展
    """
    text = Path(file_path).read_text(encoding="utf-8")
    return {
        "full_text": normalize_text(text),
        "pages": [],
        "sections": [],
        "metadata": {},
    }

def parse_pdf(file_path: str) -> dict:
    """
    解析 PDF，并保留逐页信息。

    这一步对 RAG 很关键，因为后面 chunking 可以把 page 信息带到 chunk 上，
    最终引用时就能告诉用户答案来自“第几页”。
    """
    if PdfReader is None:
        raise RuntimeError("缺少依赖 pypdf，无法解析 PDF 文件，请安装 pypdf 库")

    reader = PdfReader(file_path)
    pages: list[dict] = []
    all_parts: list[str] = []
    # 遍历每一页，提取文本并添加到 pages 列表中，enumerate是python中的内置函数，用于迭代一个序列中的元素，同时返回元素的索引和值
    for index, page in enumerate(reader.pages, start=1):
        page_text = normalize_text(page.extract_text() or "")
        if not page_text:
            continue

        pages.append(
            {
                "page": index,
                "text": page_text,
                "section": f"page_{index}",
            }
        )
        all_parts.append(page_text)

    return {
        "full_text": "\n\n".join(all_parts).strip(),
        "pages": pages,
        "sections": [],
        "metadata": {
            "parser": "pdf",
            "page_count": len(reader.pages),
        },
    }


def parse_docx(file_path: str) -> dict:
    """
    解析 DOCX，并尽量提取“章节”信息。

    简化规则：
    - 如果段落样式是 Heading*，就把它当作章节标题
    - 标题后面的正文会归到该章节下
    - 如果文档没有标题样式，就把整篇作为一个 section
    """
    if DocxDocument is None:
        raise RuntimeError("缺少依赖 python-docx，无法解析 DOCX 文件")

    document = DocxDocument(file_path)

    sections: list[dict] = []
    current_title = "正文"
    current_lines: list[str] = []

    for paragraph in document.paragraphs:
        text = normalize_text(paragraph.text or "")
        if not text:
            continue

        style_name = ""
        if paragraph.style is not None and paragraph.style.name:
            style_name = paragraph.style.name.lower()

        is_heading = style_name.startswith("heading")

        if is_heading:
            # 遇到新标题时，先把上一个 section 收起来
            if current_lines:
                section_text = normalize_text("\n".join(current_lines))
                if section_text:
                    sections.append(
                        {
                            "section": current_title,
                            "text": section_text,
                        }
                    )
            current_title = text
            current_lines = []
            continue

        current_lines.append(text)

    # 循环结束后，别忘了把最后一个 section 收尾
    if current_lines:
        section_text = normalize_text("\n".join(current_lines))
        if section_text:
            sections.append(
                {
                    "section": current_title,
                    "text": section_text,
                }
            )

    # 如果一篇文档完全没有识别出 section，就把全文作为一个 section
    if not sections:
        all_text = normalize_text("\n".join(p.text for p in document.paragraphs if (p.text or "").strip()))
        if all_text:
            sections.append(
                {
                    "section": "正文",
                    "text": all_text,
                }
            )

    full_text_parts = []
    for item in sections:
        full_text_parts.append(f"[章节] {item['section']}\n{item['text']}")

    return {
        "full_text": normalize_text("\n\n".join(full_text_parts)),
        "pages": [],
        "sections": sections,
        "metadata": {
            "parser": "docx",
            "section_count": len(sections),
        },
    }


def parse_xlsx(file_path: str) -> dict:
    """
    解析 xlsx，并把每个工作表当作一个 section。

    为什么这样设计：
    - Excel 通常没有“自然段”，最稳定的来源单位通常是“工作表”
    - 后续引用时显示 sheet 名称，用户更容易定位
    """
    if load_workbook is None:
        raise RuntimeError("缺少依赖 openpyxl，无法解析 XLSX 文件")

    workbook = load_workbook(file_path, data_only=True)
    sections: list[dict] = []

    for sheet in workbook.worksheets:
        row_lines: list[str] = []

        for row in sheet.iter_rows(values_only=True):
            cells = [_sheet_cell_to_text(value) for value in row]
            cells = [cell for cell in cells if cell]
            if not cells:
                continue

            # 用分隔符把一行单元格拼成可检索文本
            row_lines.append(" | ".join(cells))

        sheet_text = normalize_text("\n".join(row_lines))
        if not sheet_text:
            continue

        sections.append(
            {
                "section": sheet.title,
                "text": sheet_text,
            }
        )

    full_text_parts = []
    for item in sections:
        full_text_parts.append(f"[工作表] {item['section']}\n{item['text']}")

    return {
        "full_text": normalize_text("\n\n".join(full_text_parts)),
        "pages": [],
        "sections": sections,
        "metadata": {
            "parser": "xlsx",
            "sheet_count": len(sections),
        },
    }


def parse_xls(file_path: str) -> dict:
    """
    解析老版 xls。

    处理思路和 xlsx 一样：
    - 每个 sheet 作为一个 section
    - 每一行拼成一行文本
    """
    if xlrd is None:
        raise RuntimeError("缺少依赖 xlrd，无法解析 XLS 文件")

    workbook = xlrd.open_workbook(file_path)
    sections: list[dict] = []

    for sheet in workbook.sheets():
        row_lines: list[str] = []

        for row_index in range(sheet.nrows):
            values = sheet.row_values(row_index)
            cells = [_sheet_cell_to_text(value) for value in values]
            cells = [cell for cell in cells if cell]
            if not cells:
                continue

            row_lines.append(" | ".join(cells))

        sheet_text = normalize_text("\n".join(row_lines))
        if not sheet_text:
            continue

        sections.append(
            {
                "section": sheet.name,
                "text": sheet_text,
            }
        )

    full_text_parts = []
    for item in sections:
        full_text_parts.append(f"[工作表] {item['section']}\n{item['text']}")

    return {
        "full_text": normalize_text("\n\n".join(full_text_parts)),
        "pages": [],
        "sections": sections,
        "metadata": {
            "parser": "xls",
            "sheet_count": len(sections),
        },
    }


def parse_pptx(file_path: str) -> dict:
    """
    解析 PPTX，并把每一页 slide 当作一个 page。

    这样后续引用时就能输出：
    - 来自第几页 slide
    - 或者 source_section=slide_1 / slide_2
    """
    if Presentation is None:
        raise RuntimeError("缺少依赖 python-pptx，无法解析 PPTX 文件")

    presentation = Presentation(file_path)
    pages: list[dict] = []
    all_parts: list[str] = []

    for index, slide in enumerate(presentation.slides, start=1):
        slide_parts: list[str] = []

        for shape in slide.shapes:
            text = normalize_text(getattr(shape, "text", "") or "")
            if text:
                slide_parts.append(text)

        slide_text = normalize_text("\n".join(slide_parts))
        if not slide_text:
            continue

        pages.append(
            {
                "page": index,
                "text": slide_text,
                "section": f"slide_{index}",
            }
        )
        all_parts.append(slide_text)

    return {
        "full_text": normalize_text("\n\n".join(all_parts)),
        "pages": pages,
        "sections": [],
        "metadata": {
            "parser": "pptx",
            "slide_count": len(presentation.slides),
        },
    }

def parse_document(file_path: str, file_type: str) -> dict:
    """
    统一文档解析入口。

    返回结构统一为：
    {
        "full_text": str,
        "pages": list[dict],
        "sections": list[dict],
        "metadata": dict,
    }

    这样上层上传流程不用关心不同文件类型的差异，
    只要拿到统一结构后继续 chunking / embedding 即可。
    """
    file_type = (file_type or "").lower().strip()

    if file_type == "txt":
        return parse_txt(file_path)
    if file_type == "md":
        return parse_md(file_path)
    if file_type == "pdf":
        return parse_pdf(file_path)
    if file_type == "docx":
        return parse_docx(file_path)
    if file_type == "xlsx":
        return parse_xlsx(file_path)
    if file_type == "xls":
        return parse_xls(file_path)
    if file_type == "pptx":
        return parse_pptx(file_path)

    # 老版 .ppt 这里先明确报错，而不是“接口允许上传但后端悄悄失败”
    raise ValueError(f"Unsupported file type: {file_type}")